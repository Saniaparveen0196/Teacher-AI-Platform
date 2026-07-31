# app/parsers/pptx_parser.py
"""
Stage 1 — Document Intelligence: PPTX parser.
Each slide is naturally one section; slide title becomes the heading.
"""
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE


def parse_pptx(file_path: str) -> dict:
    prs = Presentation(file_path)
    sections = []
    images_or_figures = []
    full_text_parts = []

    for i, slide in enumerate(prs.slides, start=1):
        title_shape = slide.shapes.title  # may be None if layout has no title placeholder
        title = None
        body_lines = []

        for shape in slide.shapes:
            if shape.has_text_frame:
                is_title_shape = (title_shape is not None) and (shape.shape_id == title_shape.shape_id)
                for para in shape.text_frame.paragraphs:
                    line = "".join(run.text for run in para.runs).strip()
                    if not line:
                        continue
                    if is_title_shape and title is None:
                        title = line
                    else:
                        body_lines.append(line)

            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                images_or_figures.append({"page": i, "caption": None, "count": 1})

        heading = title or f"Slide {i}"
        text = "\n".join(body_lines)
        sections.append({"heading": heading, "level": 1, "text": text})
        full_text_parts.append(f"{heading}\n{text}")

    raw_text = "\n\n".join(full_text_parts)
    return {
        "raw_text": raw_text,
        "sections": sections,
        "tables": [],  # v1 skips PPTX tables — rare enough to defer
        "images_or_figures": images_or_figures,
        "metadata": {
            "page_count": len(prs.slides),
            "word_count": len(raw_text.split()),
            "source_format": "pptx",
        },
    }